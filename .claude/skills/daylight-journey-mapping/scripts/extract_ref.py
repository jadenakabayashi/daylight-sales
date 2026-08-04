#!/usr/bin/env python3
"""Multi-format text/image extraction for refs/Foundation and refs/Stakeholder
Interviews. Handles .pdf, .pptx, .docx, .xlsx — the formats that actually
show up in this project's refs/ folder. Not a general-purpose converter;
just enough to get text and embedded screenshots out reliably.

Usage:
  extract_ref.py text  <file>                     print all extractable text, page/slide/sheet-delimited
  extract_ref.py images <file> <outdir>            dump embedded images (screenshots) to outdir as PNG/JPEG
  extract_ref.py pdf-render <file> <outdir> [dpi]  render each PDF page to a PNG (for PDFs whose content
                                                    is itself one big image per page — see daylight-design-system's
                                                    experience: a "screenshot-heavy" PDF can have far fewer real
                                                    pages than the page count a naive reader assumes)

Requires: pymupdf (pdf), python-pptx (pptx), python-docx (docx), openpyxl (xlsx).
Install once with: python3 -m pip install pymupdf python-pptx python-docx openpyxl
"""
import os
import sys


def ext_of(path):
    return os.path.splitext(path)[1].lower()


def text_pdf(path):
    import fitz
    d = fitz.open(path)
    out = []
    for i, page in enumerate(d):
        out.append(f"--- page {i + 1} ---")
        out.append(page.get_text())
    return "\n".join(out)


def text_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    out = []
    for i, slide in enumerate(prs.slides):
        out.append(f"--- slide {i + 1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        out.append(line)
            if shape.has_table:
                for row in shape.table.rows:
                    out.append(" | ".join(cell.text for cell in row.cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            out.append("[notes] " + slide.notes_slide.notes_text_frame.text)
    return "\n".join(out)


def text_docx(path):
    import docx
    d = docx.Document(path)
    out = []
    for para in d.paragraphs:
        if para.text.strip():
            out.append(para.text)
    for i, table in enumerate(d.tables):
        out.append(f"--- table {i + 1} ---")
        for row in table.rows:
            out.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(out)


def text_xlsx(path):
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    out = []
    for name in wb.sheetnames:
        sheet = wb[name]
        out.append(f"--- sheet: {name} ---")
        for row in sheet.iter_rows(values_only=True):
            if any(c is not None for c in row):
                out.append("\t".join("" if c is None else str(c) for c in row))
    return "\n".join(out)


TEXT_HANDLERS = {".pdf": text_pdf, ".pptx": text_pptx, ".docx": text_docx, ".xlsx": text_xlsx}


def images_pdf(path, outdir):
    import fitz
    d = fitz.open(path)
    saved = []
    for pi, page in enumerate(d):
        for ii, img in enumerate(page.get_images(full=True)):
            xref = img[0]
            for ri, rect in enumerate(page.get_image_rects(xref)):
                pix = page.get_pixmap(matrix=fitz.Matrix(3, 3), clip=rect)
                p = os.path.join(outdir, f"p{pi + 1}_img{ii + 1}_{ri + 1}.png")
                pix.save(p)
                saved.append(p)
    return saved


def images_pptx(path, outdir):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx import Presentation
    prs = Presentation(path)
    saved = []
    for si, slide in enumerate(prs.slides):
        for shi, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext
                p = os.path.join(outdir, f"slide{si + 1}_shape{shi + 1}.{ext}")
                with open(p, "wb") as f:
                    f.write(image.blob)
                saved.append(p)
    return saved


def images_docx(path, outdir):
    import docx
    d = docx.Document(path)
    saved = []
    for i, rel in enumerate(d.part.rels.values()):
        if "image" in rel.reltype:
            blob = rel.target_part.blob
            ext = rel.target_part.content_type.split("/")[-1]
            ext = "jpg" if ext == "jpeg" else ext
            p = os.path.join(outdir, f"image{i + 1}.{ext}")
            with open(p, "wb") as f:
                f.write(blob)
            saved.append(p)
    return saved


IMAGE_HANDLERS = {".pdf": images_pdf, ".pptx": images_pptx, ".docx": images_docx}


def main():
    if len(sys.argv) < 3:
        print(__doc__, file=sys.stderr)
        sys.exit(1)
    cmd, path = sys.argv[1], sys.argv[2]
    ext = ext_of(path)

    if cmd == "text":
        handler = TEXT_HANDLERS.get(ext)
        if not handler:
            print(f"no text handler for {ext}", file=sys.stderr)
            sys.exit(1)
        print(handler(path))
    elif cmd == "images":
        if len(sys.argv) < 4:
            print("usage: extract_ref.py images <file> <outdir>", file=sys.stderr)
            sys.exit(1)
        outdir = sys.argv[3]
        os.makedirs(outdir, exist_ok=True)
        handler = IMAGE_HANDLERS.get(ext)
        if not handler:
            print(f"no image handler for {ext} (xlsx rarely embeds screenshots worth extracting)", file=sys.stderr)
            sys.exit(1)
        saved = handler(path, outdir)
        for p in saved:
            print(p)
    elif cmd == "pdf-render":
        if len(sys.argv) < 4:
            print("usage: extract_ref.py pdf-render <file> <outdir> [zoom]", file=sys.stderr)
            sys.exit(1)
        import fitz
        outdir = sys.argv[3]
        zoom = float(sys.argv[4]) if len(sys.argv) > 4 else 2.0
        os.makedirs(outdir, exist_ok=True)
        d = fitz.open(path)
        for i, page in enumerate(d):
            pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
            p = os.path.join(outdir, f"page{i + 1}.png")
            pix.save(p)
            print(p)
    else:
        print(__doc__, file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
